# -*- coding: utf-8 -*-
"""5-powerpoint.ipynb

Automatically generated by Colaboratory.

Original file is located at
    https://colab.research.google.com/github/JSJeong-me/RPA-2023/blob/main/pywinauto/5-powerpoint.ipynb
"""

# !pip install python-pptx

import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[5]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[0]

title_shape.text = '2023년도 8월 14일 전력 수요 분석 결과'

img_path = 'C:\\Users\\user\\Desktop\\result.png'

left = Inches(1)
top = Inches(1.5)

width = Inches(7)
height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)

prs.save('C:\\Users\\user\\Desktop\\new.pptxnew.pptx')